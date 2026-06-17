import io
from copy import deepcopy
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Inches, Pt

SAMPLE_TEMPLATE_PATH = Path(__file__).with_name("sample_nametag_template.pptx")
PLACEHOLDER_COMPANY = "{{company}}"
PLACEHOLDER_DEPARTMENT = "{{department}}"
PLACEHOLDER_NAME = "{{name}}"
PLACEHOLDER_DEPT_NAME = "{{dept_name}}"

Attendee = Tuple[str, str, str]


def normalize_cell(value: object) -> str:
    """Return a clean string representation for an Excel cell."""
    if pd.isna(value):
        return ""
    return str(value).strip()


def build_attendees(
    df: pd.DataFrame,
    company_col: str,
    dept_col: str,
    name_col: str,
) -> List[Attendee]:
    """Extract (company, department, name) tuples from the spreadsheet."""
    attendees: List[Attendee] = []
    for _, row in df.iterrows():
        company = normalize_cell(row[company_col])
        department = normalize_cell(row[dept_col])
        name = normalize_cell(row[name_col])
        if company and department and name:
            attendees.append((company, department, name))
    return attendees


def create_sample_template_bytes() -> bytes:
    """Create a sample nametag template PPTX with placeholders."""
    prs = Presentation()
    prs.slide_width = Inches(3.5)
    prs.slide_height = Inches(2.1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 239, 228)
    bg.line.color.rgb = RGBColor(149, 117, 59)
    bg.line.width = Pt(2)

    header = slide.shapes.add_textbox(
        Inches(0.2), Inches(0.18), Inches(3.1), Inches(0.45)
    )
    header_tf = header.text_frame
    header_tf.clear()
    header_p = header_tf.paragraphs[0]
    header_p.text = PLACEHOLDER_COMPANY
    header_p.alignment = PP_ALIGN.CENTER
    header_p.font.size = Pt(20)
    header_p.font.bold = True
    header_p.font.name = "Arial"

    title = slide.shapes.add_textbox(
        Inches(0.15), Inches(0.8), Inches(3.2), Inches(0.9)
    )
    title_tf = title.text_frame
    title_tf.clear()
    title_p = title_tf.paragraphs[0]
    title_p.text = PLACEHOLDER_DEPT_NAME
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(24)
    title_p.font.bold = True
    title_p.font.name = "Arial"

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def load_sample_template_bytes() -> bytes:
    """Load the bundled sample template, creating it in memory if missing."""
    if SAMPLE_TEMPLATE_PATH.exists():
        return SAMPLE_TEMPLATE_PATH.read_bytes()
    return create_sample_template_bytes()


def replace_placeholders_in_text(text: str, replacements: Dict[str, str]) -> str:
    """Replace template placeholders in a text run."""
    updated = text
    for placeholder, value in replacements.items():
        updated = updated.replace(placeholder, value)
    return updated


def replace_placeholders_in_shape(shape, replacements: Dict[str, str]) -> None:
    """Walk a shape tree and replace placeholders in text frames."""
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            replace_placeholders_in_shape(subshape, replacements)
        return

    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.text = replace_placeholders_in_text(run.text, replacements)
        if not paragraph.runs:
            paragraph.text = replace_placeholders_in_text(paragraph.text, replacements)


def scale_and_translate_shape(shape, scale_factor: float, dx: int, dy: int) -> None:
    """Scale a shape and its fonts, and translate it by dx, dy."""
    if hasattr(shape, "left") and shape.left is not None:
        shape.left = int(shape.left * scale_factor) + dx
    if hasattr(shape, "top") and shape.top is not None:
        shape.top = int(shape.top * scale_factor) + dy
    if hasattr(shape, "width") and shape.width is not None:
        shape.width = int(shape.width * scale_factor)
    if hasattr(shape, "height") and shape.height is not None:
        shape.height = int(shape.height * scale_factor)

    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    run.font.size = int(run.font.size * scale_factor)

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            # Subshapes in a group are scaled, but their dx/dy offset is handled by the parent group's movement.
            scale_and_translate_shape(subshape, scale_factor, 0, 0)



def remove_slide(prs: Presentation, slide_index: int) -> None:
    """Remove a slide by index using python-pptx internals."""
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)
    slide_id_list.remove(slides[slide_index])


def generate_ppt_from_template(
    attendees: List[Attendee], 
    template_bytes: bytes,
    grid_rows: int,
    grid_cols: int,
    is_landscape: bool
) -> io.BytesIO:
    """Generate PPTX containing nametags arranged in an A4 grid."""
    prs = Presentation(io.BytesIO(template_bytes))
    if not prs.slides:
        raise ValueError("The template PPTX must contain at least one slide.")

    template_slide = prs.slides[0]
    template_w = prs.slide_width
    template_h = prs.slide_height

    a4_w = Cm(29.7) if is_landscape else Cm(21.0)
    a4_h = Cm(21.0) if is_landscape else Cm(29.7)

    prs.slide_width = a4_w
    prs.slide_height = a4_h

    cell_w = a4_w / grid_cols
    cell_h = a4_h / grid_rows

    scale_w = cell_w / template_w
    scale_h = cell_h / template_h
    scale_factor = min(scale_w, scale_h)

    scaled_w = template_w * scale_factor
    scaled_h = template_h * scale_factor

    nametags_per_page = grid_rows * grid_cols
    blank_layout = prs.slide_layouts[6]

    for i in range(0, len(attendees), nametags_per_page):
        chunk = attendees[i:i + nametags_per_page]
        new_slide = prs.slides.add_slide(blank_layout)

        for j, (company, department, name) in enumerate(chunk):
            row = j // grid_cols
            col = j % grid_cols

            dx = int((col * cell_w) + (cell_w - scaled_w) / 2)
            dy = int((row * cell_h) + (cell_h - scaled_h) / 2)

            start_idx = len(new_slide.shapes)
            for shape in template_slide.shapes:
                new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
            end_idx = len(new_slide.shapes)

            replacements = {
                PLACEHOLDER_COMPANY: company,
                PLACEHOLDER_DEPARTMENT: department,
                PLACEHOLDER_NAME: name,
                PLACEHOLDER_DEPT_NAME: f"{department} {name}",
            }

            for shape_idx in range(start_idx, end_idx):
                shape = new_slide.shapes[shape_idx]
                scale_and_translate_shape(shape, scale_factor, dx, dy)
                replace_placeholders_in_shape(shape, replacements)

    remove_slide(prs, 0)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def detect_default_index(columns: List[str], candidates: List[str], fallback: int = 0) -> int:
    """Return the first matching column index for the provided aliases."""
    lowered = {str(col).strip().lower(): idx for idx, col in enumerate(columns)}
    for candidate in candidates:
        match = lowered.get(candidate.lower())
        if match is not None:
            return match
    return fallback


def load_dataframe(uploaded_file) -> pd.DataFrame:
    """Read an uploaded CSV or Excel file into a DataFrame."""
    if uploaded_file.name.lower().endswith(".csv"):
        return pd.read_csv(uploaded_file)
    return pd.read_excel(uploaded_file)


def read_template_bytes(uploaded_template) -> bytes:
    """Return uploaded template bytes or the bundled sample template bytes."""
    if uploaded_template is None:
        return load_sample_template_bytes()
    return uploaded_template.getvalue()


def main() -> None:
    import streamlit as st

    st.set_page_config(page_title="Nametag Generator")
    
    st.markdown("""
    <div style="text-align: center; padding: 16px 0 20px;">
        <h1 style="color: #171717; font-family: Arial, sans-serif; margin-bottom: 0;">명찰 생성기</h1>
        <p style="color: #6b7280; margin-top: 8px;">Nametag Generator</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown(
        """
        Excel 파일의 `회사명`, `부서`, `이름` 컬럼을 읽어 명찰을 생성합니다.
        템플릿 PPTX의 첫 번째 슬라이드를 기준으로 명찰을 복제하며,
        `{{company}}`, `{{department}}`, `{{name}}`, `{{dept_name}}` 플레이스홀더를 치환합니다.
        결과는 모든 명찰이 합쳐진 하나의 PPTX 파일로 다운로드됩니다.
        """
    )

    sample_template_bytes = load_sample_template_bytes()
    st.download_button(
        label="샘플 템플릿 다운로드",
        data=sample_template_bytes,
        file_name="sample_nametag_template.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    st.caption(
        "샘플 템플릿의 맨 위는 회사명, 가운데는 부서+이름 자리입니다. "
        "직접 만든 템플릿을 쓸 때는 첫 번째 슬라이드에 플레이스홀더를 넣어주세요."
    )

    uploaded_file = st.file_uploader("명단 Excel/CSV 업로드", type=["xlsx", "xls", "csv"])
    uploaded_template = st.file_uploader("템플릿 PPTX 업로드", type=["pptx"])

    if uploaded_file:
        try:
            df = load_dataframe(uploaded_file)
        except Exception as exc:
            st.error(f"파일을 읽지 못했습니다: {exc}")
            return

        st.success(f"Loaded file with {len(df)} rows and {len(df.columns)} columns.")

        columns = df.columns.tolist()
        company_col = st.selectbox(
            "회사명 컬럼",
            columns,
            index=detect_default_index(columns, ["회사명", "company"], 0),
        )
        dept_col = st.selectbox(
            "부서 컬럼",
            columns,
            index=detect_default_index(columns, ["부서", "department", "dept"], min(1, len(columns) - 1)),
        )
        name_col = st.selectbox(
            "이름 컬럼",
            columns,
            index=detect_default_index(columns, ["이름", "name"], min(2, len(columns) - 1)),
        )

        preview_df = df[[company_col, dept_col, name_col]].rename(
            columns={company_col: "회사명", dept_col: "부서", name_col: "이름"}
        )
        st.dataframe(preview_df.head(10), use_container_width=True)

        col1, col2, col3 = st.columns(3)
        with col1:
            orientation = st.radio("용지 방향", ["가로 (Landscape)", "세로 (Portrait)"], index=0)
        with col2:
            grid_size = st.selectbox("명찰 배열 (가로칸 x 세로칸)", ["1x1", "1x2", "2x1", "1x3", "3x1", "1x4", "4x1", "2x2", "2x3", "3x2", "2x4", "4x2", "3x3"])
        with col3:
            num_copies = st.number_input("인당 출력 개수", min_value=1, max_value=10, value=1)
        
        grid_cols, grid_rows = map(int, grid_size.split("x"))
        is_landscape = orientation.startswith("가로")

        attendees = build_attendees(df, company_col, dept_col, name_col)
        # 인당 출력 개수만큼 복제
        attendees = [a for a in attendees for _ in range(num_copies)]

        # 미리보기 로직
        st.subheader("예상 도안 미리보기 (A4 용지 기준)")
        template_bytes = read_template_bytes(uploaded_template)
        try:
            prs_preview = Presentation(io.BytesIO(template_bytes))
            template_w = prs_preview.slide_width
            template_h = prs_preview.slide_height
        except Exception:
            template_w = Inches(3.5)
            template_h = Inches(2.1)

        EMU_PER_CM = 360000
        A4_W_EMU = int(29.7 * EMU_PER_CM) if is_landscape else int(21.0 * EMU_PER_CM)
        A4_H_EMU = int(21.0 * EMU_PER_CM) if is_landscape else int(29.7 * EMU_PER_CM)

        cell_w = A4_W_EMU / grid_cols
        cell_h = A4_H_EMU / grid_rows

        scale_w = cell_w / template_w
        scale_h = cell_h / template_h
        scale_factor = min(scale_w, scale_h)

        scaled_w = template_w * scale_factor
        scaled_h = template_h * scale_factor

        scaled_w_cm = scaled_w / EMU_PER_CM
        scaled_h_cm = scaled_h / EMU_PER_CM

        st.info(f"선택한 배열에 따른 개별 명찰 출력 크기: **가로 {scaled_w_cm:.1f}cm x 세로 {scaled_h_cm:.1f}cm**")

        display_h = 350
        display_w = display_h * (A4_W_EMU / A4_H_EMU)

        import textwrap
        
        html_content = textwrap.dedent(f'''
        <div style="
            width: {display_w}px;
            height: {display_h}px;
            border: 2px solid #ccc;
            background-color: #fafafa;
            margin: 0 auto;
            margin-bottom: 20px;
            position: relative;
            box-shadow: 2px 2px 10px rgba(0,0,0,0.1);
        ">
        ''')

        cell_disp_w = display_w / grid_cols
        cell_disp_h = display_h / grid_rows

        nametag_disp_w = cell_disp_w * (scaled_w / cell_w)
        nametag_disp_h = cell_disp_h * (scaled_h / cell_h)

        preview_index = 0
        for r in range(grid_rows):
            for c in range(grid_cols):
                dx = (c * cell_disp_w) + (cell_disp_w - nametag_disp_w) / 2
                dy = (r * cell_disp_h) + (cell_disp_h - nametag_disp_h) / 2
                base_font_size = max(8, int(nametag_disp_h * 0.12))
                span_font_size = max(10, int(nametag_disp_h * 0.18))
                
                # 명찰에 들어갈 실제 텍스트 가져오기
                if preview_index < len(attendees):
                    comp, dept, name = attendees[preview_index]
                    display_text = f"{comp}<br><span style='font-size:{span_font_size}px; color:#333;'>{dept} {name}</span>"
                else:
                    display_text = ""
                
                html_content += textwrap.dedent(f'''
                <div style="
                    position: absolute;
                    left: {dx}px;
                    top: {dy}px;
                    width: {nametag_disp_w}px;
                    height: {nametag_disp_h}px;
                    background-color: #ffffff;
                    border: 1px dashed #9ca3af;
                    border-radius: 4px;
                    display: flex;
                    align-items: center;
                    justify-content: center;
                    flex-direction: column;
                    text-align: center;
                    font-size: {base_font_size}px;
                    font-weight: bold;
                    color: #111827;
                    padding: 10px;
                    box-sizing: border-box;
                    overflow: hidden;
                ">
                    {display_text}
                </div>
                ''')
                preview_index += 1
        html_content += "</div>"
        
        import streamlit.components.v1 as components
        components.html(html_content, height=int(display_h) + 30)

        if st.button("명찰 만들기 (다운로드 파일 생성)", type="primary"):
            if not attendees:
                st.warning("회사명, 부서, 이름이 모두 채워진 행을 찾지 못했습니다.")
                return

            try:
                result = generate_ppt_from_template(
                    attendees, 
                    template_bytes, 
                    grid_rows=grid_rows, 
                    grid_cols=grid_cols, 
                    is_landscape=is_landscape
                )
            except Exception as exc:
                st.error(f"명찰 생성에 실패했습니다: {exc}")
                return

            st.success(f"{len(attendees)}개의 명찰을 하나의 PPTX 파일로 생성했습니다.")
            st.download_button(
                label="명찰 PPTX 다운로드",
                data=result,
                file_name="nametags.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )


if __name__ == "__main__":
    main()
